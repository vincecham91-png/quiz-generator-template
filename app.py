import os
import re
import random
from flask import Flask, request, render_template, redirect, url_for, flash
from werkzeug.utils import secure_filename

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXT = {'.pdf', '.pptx', '.docx', '.doc'}

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.secret_key = 'change-this-in-production'

os.makedirs(UPLOAD_FOLDER, exist_ok=True)


def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    text = ''
    if ext == '.pdf':
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            for p in reader.pages:
                text += p.extract_text() or '\n'
        except Exception as e:
            text = ''
    elif ext == '.docx':
        try:
            from docx import Document
            doc = Document(path)
            for p in doc.paragraphs:
                text += p.text + '\n'
        except Exception:
            text = ''
    elif ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
        except Exception:
            text = ''
    elif ext == '.doc':
        # .doc (legacy) fallback: ask user to convert to .docx for best results
        text = ''

    return text


def make_questions_from_text(text, max_q=10):
    if not text or len(text.strip()) < 50:
        return []

    # split into sentences (handles Chinese punctuation and Latin punctuation)
    sents = re.split(r'(?<=[。！？!?\.])\s*', text)
    candidates = [s.strip() for s in sents if len(s.strip()) >= 10]
    questions = []

    # collect Chinese word-like fragments for options pool
    pool = re.findall(r'[\u4e00-\u9fff]{2,6}', text)
    # fallback for Latin words
    pool += re.findall(r"[A-Za-z]{4,}", text)

    random.shuffle(pool)

    for sent in candidates:
        if len(questions) >= max_q:
            break
        # try to find a Chinese fragment inside sentence
        m = re.findall(r'[\u4e00-\u9fff]{2,6}', sent)
        answer = None
        if m:
            answer = random.choice(m)
        else:
            # fallback to English word
            m2 = re.findall(r"[A-Za-z]{4,}", sent)
            if m2:
                answer = random.choice(m2)
        if not answer:
            continue

        # create options: correct + 3 distractors from pool
        distractors = [w for w in pool if w != answer]
        random.shuffle(distractors)
        opts = [answer]
        for d in distractors[:3]:
            opts.append(d)
        if len(opts) < 4:
            continue
        random.shuffle(opts)
        question_text = sent.replace(answer, '____')
        q = {
            'question': question_text,
            'options': opts,
            'answer': opts.index(answer),
            'explanation': f'原句：{sent}'
        }
        questions.append(q)

    return questions


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        flash('No file part')
        return redirect(url_for('index'))
    f = request.files['file']
    if f.filename == '':
        flash('No selected file')
        return redirect(url_for('index'))
    filename = secure_filename(f.filename)
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_EXT:
        flash('File type not allowed')
        return redirect(url_for('index'))
    save_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    f.save(save_path)

    text = extract_text(save_path)
    questions = make_questions_from_text(text)
    if not questions:
        flash('Could not extract enough text. Try a different file or convert .doc to .docx')
        return redirect(url_for('index'))

    return render_template('quiz.html', questions=questions, title=filename)


if __name__ == '__main__':
    app.run(debug=True)
